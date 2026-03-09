"""
generate_ppt.py — 雲智判上線自動化流程魚骨圖報告 PPT 生成腳本
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 顏色定義 ─────────────────────────────────────────────────────────────────
C_DARK    = RGBColor(0x2C, 0x3E, 0x50)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xF0, 0xF4, 0xF8)
C_PURPLE  = RGBColor(0x8E, 0x44, 0xAD)
C_RED     = RGBColor(0xE7, 0x4C, 0x3C)
C_ORANGE  = RGBColor(0xE6, 0x7E, 0x22)
C_YELLOW  = RGBColor(0xD4, 0xAC, 0x0D)
C_GREEN   = RGBColor(0x27, 0xAE, 0x60)
C_BLUE    = RGBColor(0x29, 0x80, 0xB9)
C_GRAY    = RGBColor(0x95, 0xA5, 0xA6)

W = Inches(13.33)
H = Inches(7.5)

STAGES = [
    {
        "no": "01", "name": "初期評估", "color": C_RED,
        "current": "由開發與規劃人員以人工方式進行需求評估",
        "problems": ["缺乏量化評估指標", "缺乏模型可行性評估工具", "完全依賴人工主觀判斷", "評估時程冗長"],
        "missing":  ["自動化指標評估模組（DINOv3）", "數據可行性評分系統"],
    },
    {
        "no": "02", "name": "數據標註", "color": C_ORANGE,
        "current": "評估完成後由開發人員全程人工標註",
        "problems": ["完全依賴人工手動標註", "缺乏通用標註方法 / 框架", "標註一致性難以保證", "標註效率低，開發時程長"],
        "missing":  ["半自動 / 主動學習標註工具", "通用標註標準與流程"],
    },
    {
        "no": "03", "name": "模型訓練", "color": C_YELLOW,
        "current": "前期需要人工撰寫訓練程式碼，超參數手動調整",
        "problems": ["需要手動撰寫訓練程式碼", "超參數調整依賴人工經驗", "缺乏自動化訓練流水線", "重複性工作佔用大量開發資源"],
        "missing":  ["AutoML / 訓練流水線自動化", "超參數自動搜索（HPO）"],
    },
    {
        "no": "04", "name": "模型優化", "color": C_GREEN,
        "current": "上線後設備差異造成大量過殺，需人工不斷重新 Label",
        "problems": ["設備差異導致大量過殺", "曝光度 / 機台調整造成圖像偏移", "新機台須大量重新人工 Label", "10 台機台只標 2 台，其餘補標成本高", "缺乏自動化優化與輔助標註方案"],
        "missing":  ["自動域適應（Domain Adaptation）", "AI 輔助標註 / 主動學習", "多機台差異自動補償"],
    },
    {
        "no": "05", "name": "用戶反饋", "color": C_BLUE,
        "current": "上線後僅能取得「復判 OK」資訊，漏檢 Sorting 無模型支援",
        "problems": ["僅能獲得「復判 OK」反饋資訊", "缺乏漏檢 Sorting 機制", "過殺 / 漏檢反饋無法閉環進模型", "缺乏相似異常物料自動搜尋能力"],
        "missing":  ["OK/NG 雙向反饋閉環系統", "漏檢異常 Sorting 模型", "持續學習（Continual Learning）"],
    },
]


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None,
             line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=Pt(14), bold=False, color=None,
                 align=PP_ALIGN.LEFT, font_name="微軟正黑體"):
    if color is None:
        color = C_DARK
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = font_size
    run.font.bold      = bold
    run.font.color.rgb = color
    run.font.name      = font_name
    return txb


def add_label_box(slide, text, left, top, width, height,
                  bg_color, font_size=Pt(13), font_name="微軟正黑體"):
    rect = add_rect(slide, left, top, width, height, fill_color=bg_color)
    tf = rect.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text           = text
    run.font.size      = font_size
    run.font.bold      = True
    run.font.color.rgb = C_WHITE
    run.font.name      = font_name
    return rect


prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]

# ── Slide 1: 封面 ─────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, C_DARK)
add_rect(s1, 0, Inches(2.6), W, Pt(4), fill_color=C_PURPLE)
add_text_box(s1, "雲智判上線自動化流程",
             Inches(1.5), Inches(1.5), Inches(10), Inches(1),
             font_size=Pt(40), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(s1, "魚骨圖分析 — 各階段缺失問題與建議補足技術",
             Inches(2), Inches(2.8), Inches(9.3), Inches(0.8),
             font_size=Pt(20), color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)
stage_colors = [C_RED, C_ORANGE, C_YELLOW, C_GREEN, C_BLUE]
stage_names  = ["初期評估", "數據標註", "模型訓練", "模型優化", "用戶反饋"]
for i, (c, n) in enumerate(zip(stage_colors, stage_names)):
    add_label_box(s1, n, Inches(0.6) + i * Inches(2.42), Inches(3.9),
                  Inches(2.2), Inches(0.55), bg_color=c, font_size=Pt(13))
add_text_box(s1, "2026", Inches(0), Inches(6.8), W, Inches(0.5),
             font_size=Pt(13), color=C_GRAY, align=PP_ALIGN.CENTER)

# ── Slide 2: 流程概覽 ─────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, C_LIGHT)
add_rect(s2, 0, 0, W, Inches(1.1), fill_color=C_DARK)
add_text_box(s2, "自動化流程五大階段",
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.8),
             font_size=Pt(28), bold=True, color=C_WHITE)
BOX_W = Inches(2.0)
BOX_H = Inches(1.5)
GAP   = Inches(0.3)
TOP_Y = Inches(2.0)
descriptions = [
    "人工需求評估\n缺乏量化指標",
    "人工手動標註\n缺通用框架",
    "人工撰寫訓練\n缺自動化流水線",
    "設備差異造成過殺\n缺輔助標註方案",
    "反饋資訊不完整\n缺漏檢Sorting",
]
for i, (c, n, d) in enumerate(zip(stage_colors, stage_names, descriptions)):
    x = Inches(0.5) + i * (BOX_W + GAP)
    rect = add_rect(s2, x, TOP_Y, BOX_W, BOX_H, fill_color=c)
    tf = rect.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text           = f"{i+1:02d}. {n}"
    run.font.size      = Pt(14)
    run.font.bold      = True
    run.font.color.rgb = C_WHITE
    run.font.name      = "微軟正黑體"
    add_text_box(s2, d, x, TOP_Y + BOX_H + Inches(0.1), BOX_W, Inches(0.8),
                 font_size=Pt(10.5), align=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(s2, "▶", x + BOX_W + Inches(0.05),
                     TOP_Y + BOX_H / 2 - Inches(0.2), GAP, Inches(0.4),
                     font_size=Pt(14), align=PP_ALIGN.CENTER)
add_text_box(s2, "✗ 現存問題　　→ 建議補足技術",
             Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
             font_size=Pt(12), color=C_GRAY, align=PP_ALIGN.CENTER)

# ── Slides 3-7: 各階段詳細頁 ──────────────────────────────────────────────────
for st in STAGES:
    s = prs.slides.add_slide(blank_layout)
    set_bg(s, C_LIGHT)
    add_rect(s, 0, 0, W, Inches(1.2), fill_color=st["color"])
    add_text_box(s, f"{st['no']}  {st['name']}",
                 Inches(0.5), Inches(0.2), Inches(8), Inches(0.8),
                 font_size=Pt(28), bold=True, color=C_WHITE)
    add_text_box(s, f"現況：{st['current']}",
                 Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.5),
                 font_size=Pt(12))
    add_rect(s, Inches(0.5), Inches(1.82), Inches(12.3), Pt(2),
             fill_color=st["color"])

    # 左欄：現存問題
    add_text_box(s, "現存問題",
                 Inches(0.5), Inches(1.95), Inches(6), Inches(0.4),
                 font_size=Pt(14), bold=True, color=st["color"])
    for i, prob in enumerate(st["problems"]):
        py = Inches(2.45) + i * Inches(0.8)
        icon = add_rect(s, Inches(0.5), py, Inches(0.4), Inches(0.55),
                        fill_color=st["color"])
        ti = icon.text_frame.paragraphs[0]
        ti.alignment = PP_ALIGN.CENTER
        ri = ti.add_run()
        ri.text           = "✗"
        ri.font.size      = Pt(14)
        ri.font.bold      = True
        ri.font.color.rgb = C_WHITE
        ri.font.name      = "微軟正黑體"
        add_text_box(s, prob, Inches(1.0), py + Inches(0.05),
                     Inches(5.7), Inches(0.55), font_size=Pt(12.5))

    # 右欄：建議補足技術
    add_rect(s, Inches(7.2), Inches(1.95), Inches(5.6), Inches(0.4),
             fill_color=st["color"])
    add_text_box(s, "建議補足技術",
                 Inches(7.2), Inches(1.95), Inches(5.6), Inches(0.4),
                 font_size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    for i, miss in enumerate(st["missing"]):
        my = Inches(2.55) + i * Inches(1.1)
        card = add_rect(s, Inches(7.2), my, Inches(5.6), Inches(0.9),
                        fill_color=C_WHITE, line_color=st["color"], line_width=Pt(1.5))
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text           = f"→  {miss}"
        run.font.size      = Pt(12)
        run.font.bold      = True
        run.font.color.rgb = st["color"]
        run.font.name      = "微軟正黑體"

# ── Slide 8: 魚骨圖 ───────────────────────────────────────────────────────────
import os
s8 = prs.slides.add_slide(blank_layout)
set_bg(s8, C_DARK)
add_rect(s8, 0, 0, W, Inches(0.9), fill_color=C_PURPLE)
add_text_box(s8, "魚骨圖總覽",
             Inches(0.5), Inches(0.1), Inches(12), Inches(0.7),
             font_size=Pt(26), bold=True, color=C_WHITE)
if os.path.exists("fishbone_diagram.png"):
    s8.shapes.add_picture("fishbone_diagram.png", Inches(0.3), Inches(1.0),
                           width=Inches(12.7))
else:
    add_text_box(s8, "（請先執行 fishbone_diagram.py 生成圖片後重新執行本腳本）",
                 Inches(2), Inches(3), Inches(9), Inches(1),
                 font_size=Pt(14), color=C_GRAY, align=PP_ALIGN.CENTER)

# ── Slide 9: 總結與建議 ───────────────────────────────────────────────────────
s9 = prs.slides.add_slide(blank_layout)
set_bg(s9, C_LIGHT)
add_rect(s9, 0, 0, W, Inches(1.0), fill_color=C_DARK)
add_text_box(s9, "總結與優先建議",
             Inches(0.5), Inches(0.12), Inches(12), Inches(0.75),
             font_size=Pt(26), bold=True, color=C_WHITE)
PRIORITIES = [
    (C_RED,    "初期評估",  "導入 DINOv3 自動化評估模組，量化 OK/NG 可分性與數據複雜度，縮短評估時程"),
    (C_ORANGE, "數據標註",  "建立半自動主動學習標註流程，統一標註標準，降低人工依賴"),
    (C_YELLOW, "模型訓練",  "建置 AutoML 訓練流水線，整合超參數自動搜索（HPO），減少重複開發"),
    (C_GREEN,  "模型優化",  "導入 Domain Adaptation 自動補償設備差異，AI 輔助標註解決多機台補標問題"),
    (C_BLUE,   "用戶反饋",  "建立 OK/NG 雙向反饋閉環，加入漏檢 Sorting 模型，推動持續學習機制"),
]
for i, (c, name, desc) in enumerate(PRIORITIES):
    row_y = Inches(1.15) + i * Inches(1.12)
    add_label_box(s9, name, Inches(0.4), row_y, Inches(1.5), Inches(0.75), bg_color=c, font_size=Pt(12))
    add_text_box(s9, desc, Inches(2.1), row_y + Inches(0.05), Inches(10.8), Inches(0.75), font_size=Pt(12))
    if i < 4:
        add_rect(s9, Inches(0.4), row_y + Inches(0.78), Inches(12.5), Pt(1),
                 fill_color=RGBColor(0xCC, 0xCC, 0xCC))

# ── Slide 10: 結尾頁 ──────────────────────────────────────────────────────────
s10 = prs.slides.add_slide(blank_layout)
set_bg(s10, C_DARK)
add_rect(s10, 0, Inches(3.1), W, Pt(3), fill_color=C_PURPLE)
add_text_box(s10, "Thank You",
             Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
             font_size=Pt(52), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text_box(s10, "雲智判上線自動化流程 — 持續優化，共同推進",
             Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.7),
             font_size=Pt(16), color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)
for i, (c, n) in enumerate(zip(stage_colors, stage_names)):
    add_label_box(s10, n, Inches(0.9) + i * Inches(2.42), Inches(4.5),
                  Inches(2.2), Inches(0.5), bg_color=c, font_size=Pt(12))

out = "雲智判上線自動化流程分析.pptx"
prs.save(out)
print(f"[✓] PPT 已儲存至: {out}")
